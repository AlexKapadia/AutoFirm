"""Structural + determinism tests for the slide-deck builder.

Every generated deck is RE-OPENED with python-pptx and its real structure
asserted — slide count, the title slide, each action title, each bullet's text,
the 16:9 canvas, the applied design-token font/colour, and the killed default
drop-shadow — not merely that ``build`` ran. Determinism and fail-closed paths
are exercised (CLAUDE.md §3.6, §5.6).
"""

from __future__ import annotations

from pathlib import Path

import pytest
from hypothesis import HealthCheck, given, settings
from hypothesis import strategies as st
from pptx import Presentation

from autofirm.artifacts import slide_deck_design_tokens as tokens
from autofirm.artifacts.artifact_spec_validation_errors import ArtifactSpecError
from autofirm.artifacts.slide_deck_builder import build_slide_deck
from autofirm.artifacts.slide_deck_spec import DeckSpec, SlideSpec


def _deck() -> DeckSpec:
    return DeckSpec(
        title="Acme Q4 Review",
        subtitle="Board pack — Dec 2024",
        slides=(
            SlideSpec(
                "Revenue grew thirty two percent on enterprise",
                ("Enterprise ARR up 48%", "Net retention 121%"),
            ),
            SlideSpec("Cash runway now extends to twenty six months", ("Burn cut 18%",)),
        ),
    )


def _texts(slide: object) -> list[str]:
    return [sh.text_frame.text for sh in slide.shapes if sh.has_text_frame]  # type: ignore[attr-defined]


def test_deck_has_title_slide_plus_one_per_content_slide(tmp_path: Path) -> None:
    path = build_slide_deck(_deck(), tmp_path / "d.pptx")
    pres = Presentation(path)
    assert len(pres.slides._sldIdLst) == 3  # 1 title + 2 content


def test_title_slide_carries_title_and_subtitle(tmp_path: Path) -> None:
    path = build_slide_deck(_deck(), tmp_path / "d.pptx")
    first = next(iter(Presentation(path).slides))
    texts = _texts(first)
    assert "Acme Q4 Review" in texts
    assert any("Board pack" in t for t in texts)


def test_content_slide_action_titles_present(tmp_path: Path) -> None:
    path = build_slide_deck(_deck(), tmp_path / "d.pptx")
    slides = list(Presentation(path).slides)
    assert any("Revenue grew thirty two percent" in t for t in _texts(slides[1]))
    assert any("Cash runway" in t for t in _texts(slides[2]))


def test_bullets_rendered_with_glyph(tmp_path: Path) -> None:
    path = build_slide_deck(_deck(), tmp_path / "d.pptx")
    slides = list(Presentation(path).slides)
    body = "\n".join(_texts(slides[1]))
    assert "Enterprise ARR up 48%" in body
    assert "Net retention 121%" in body
    assert "•" in body  # explicit bullet glyph from our owned layout


def test_canvas_is_16_by_9(tmp_path: Path) -> None:
    path = build_slide_deck(_deck(), tmp_path / "d.pptx")
    pres = Presentation(path)
    assert pres.slide_width == tokens.SLIDE_WIDTH
    assert pres.slide_height == tokens.SLIDE_HEIGHT


def test_title_run_uses_design_tokens(tmp_path: Path) -> None:
    path = build_slide_deck(_deck(), tmp_path / "d.pptx")
    slides = list(Presentation(path).slides)
    # Find the action-title text frame and assert the template font + ink colour.
    for shape in slides[1].shapes:
        if shape.has_text_frame and "Revenue grew" in shape.text_frame.text:
            run = shape.text_frame.paragraphs[0].runs[0]
            assert run.font.name == tokens.HEADING_FONT
            assert run.font.bold is True
            assert run.font.color.rgb == tokens.INK
            return
    pytest.fail("action-title text frame not found")


def test_no_inherited_placeholders(tmp_path: Path) -> None:
    # We build from the BLANK layout, so content slides own only the shapes we
    # add — no theme Accent placeholders bleeding in.
    path = build_slide_deck(_deck(), tmp_path / "d.pptx")
    slides = list(Presentation(path).slides)
    assert len(list(slides[1].placeholders)) == 0


def test_accent_rule_has_no_drop_shadow(tmp_path: Path) -> None:
    # Anti-slop: the default autoshape drop-shadow must be suppressed.
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    path = build_slide_deck(_deck(), tmp_path / "d.pptx")
    slides = list(Presentation(path).slides)
    rules = [s for s in slides[1].shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    assert rules, "expected the accent-rule autoshape on the content slide"
    assert rules[0].shadow.inherit is False


def test_deck_build_is_deterministic(tmp_path: Path) -> None:
    a = build_slide_deck(_deck(), tmp_path / "a.pptx")
    b = build_slide_deck(_deck(), tmp_path / "b.pptx")
    assert a.read_bytes() == b.read_bytes()


@pytest.mark.parametrize("name", ["d.txt", "d.ppt", "d", "deck.key"])
def test_non_pptx_destination_refused(tmp_path: Path, name: str) -> None:
    with pytest.raises(ArtifactSpecError, match=r"must be a \.pptx"):
        build_slide_deck(_deck(), tmp_path / name)


def test_title_only_deck_omits_subtitle_box(tmp_path: Path) -> None:
    spec = DeckSpec(title="No Subtitle Deck", slides=(SlideSpec("A valid claim here now", ("x",)),))
    path = build_slide_deck(spec, tmp_path / "d.pptx")
    first = next(iter(Presentation(path).slides))
    # Only the title text box exists; no empty subtitle frame is emitted.
    text_frames = [sh for sh in first.shapes if sh.has_text_frame]
    assert len(text_frames) == 1


@pytest.mark.property
@settings(max_examples=30, suppress_health_check=[HealthCheck.function_scoped_fixture])
@given(n_slides=st.integers(min_value=1, max_value=6))
def test_slide_count_matches_spec(tmp_path: Path, n_slides: int) -> None:
    slides = tuple(
        SlideSpec(f"This is claim number {i} here", (f"point {i}",)) for i in range(n_slides)
    )
    spec = DeckSpec(title="Generated", slides=slides)
    path = build_slide_deck(spec, tmp_path / "g.pptx")
    pres = Presentation(path)
    # 1 title slide + one per content slide, for any deck size.
    assert len(pres.slides._sldIdLst) == n_slides + 1
