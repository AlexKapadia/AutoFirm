"""Render a Minto-structured :class:`DeckSpec` to a ``.pptx`` slide deck.

What this does
--------------
Builds a real PowerPoint deck from a validated deck spec using a *custom blank
template* and explicit placement, not python-pptx's built-in bullet layouts.
Every slide leads with its action title under an accent rule, then lays MECE
bullets on a shared spacing grid (``slide_deck_design_tokens``). The result has a
deliberate type/spacing hierarchy (CLAUDE.md §3.14;
``docs/research/B15-artifact-generation`` §2.2), not library defaults.

Why it exists / where it sits
-----------------------------
This is the L1.B15.2 deck builder. It starts from the blank slide layout so it
owns every shape (no inherited theme placeholders / Accent 1-6 slop), giving
byte-deterministic output for a given spec (CLAUDE.md §3.6).

Security / compliance invariants upheld
---------------------------------------
The builder refuses a non-``.pptx`` target *before* writing (fail-closed —
CLAUDE.md §5.6). All text is written through python-pptx text frames (no template
string interpolation), so spec content cannot alter slide structure.
"""

from __future__ import annotations

from pathlib import Path
from typing import TYPE_CHECKING

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Length, Pt

from autofirm.artifacts import slide_deck_design_tokens as tokens
from autofirm.artifacts.artifact_spec_validation_errors import ArtifactSpecError

if TYPE_CHECKING:
    from pptx.slide import Slide
    from pptx.text.text import _Run

    from autofirm.artifacts.slide_deck_spec import DeckSpec, SlideSpec

__all__ = ["build_slide_deck"]

_BLANK_LAYOUT_INDEX = 6  # the built-in fully-blank layout: we own every shape


def _emu_sum(*lengths: Length) -> Length:
    """Add EMU ``Length`` values and return a ``Length`` (not a bare ``int``).

    python-pptx's ``Length`` arithmetic yields plain ``int`` EMU; geometry helpers
    require ``Length``, so we re-wrap. Keeps the layout maths typed and exact.
    """
    return Emu(sum(int(length) for length in lengths))


def build_slide_deck(spec: DeckSpec, destination: Path) -> Path:
    """Write ``spec`` to ``destination`` as a designed ``.pptx`` deck.

    Args:
        spec: A validated deck spec.
        destination: Target path; must end ``.pptx``. Parent dirs are created.

    Returns:
        The ``destination`` path (now a written deck).

    Raises:
        ArtifactSpecError: If ``destination`` is not a ``.pptx`` path (fail-closed
            — CLAUDE.md §5.6).
    """
    if destination.suffix.lower() != ".pptx":  # fail-closed: wrong container refused
        raise ArtifactSpecError(f"destination must be a .pptx file, got {destination.name!r}")

    presentation = Presentation()
    presentation.slide_width = tokens.SLIDE_WIDTH
    presentation.slide_height = tokens.SLIDE_HEIGHT
    blank_layout = presentation.slide_layouts[_BLANK_LAYOUT_INDEX]

    _render_title_slide(presentation.slides.add_slide(blank_layout), spec)
    for slide_spec in spec.slides:
        _render_content_slide(presentation.slides.add_slide(blank_layout), slide_spec)

    destination.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(destination))  # python-pptx save() takes a path string
    return destination


def _render_title_slide(slide: Slide, spec: DeckSpec) -> None:
    """Place the deck title (and optional subtitle) on a clean cover."""
    title_top = Emu(int(tokens.SLIDE_HEIGHT) // 3)  # vertically centre-ish anchor
    title_height = _emu_sum(tokens.TITLE_SIZE, tokens.TITLE_SIZE)
    title_box = slide.shapes.add_textbox(
        tokens.CONTENT_LEFT, title_top, tokens.CONTENT_WIDTH, title_height
    )
    para = title_box.text_frame.paragraphs[0]
    para.text = spec.title
    _style_run(para.runs[0], tokens.TITLE_SIZE, tokens.INK, bold=True)

    if spec.subtitle.strip():
        sub_box = slide.shapes.add_textbox(
            tokens.CONTENT_LEFT,
            _emu_sum(title_top, title_height),
            tokens.CONTENT_WIDTH,
            _emu_sum(tokens.SUBTITLE_SIZE, tokens.SUBTITLE_SIZE),
        )
        sub_para = sub_box.text_frame.paragraphs[0]
        sub_para.text = spec.subtitle
        _style_run(sub_para.runs[0], tokens.SUBTITLE_SIZE, tokens.MUTED, bold=False)


def _render_content_slide(slide: Slide, slide_spec: SlideSpec) -> None:
    """Place the action title, an accent rule, and the MECE bullets."""
    _add_action_title(slide, slide_spec.action_title)
    _add_accent_rule(slide)
    _add_bullets(slide, slide_spec.bullets)


def _add_action_title(slide: Slide, action_title: str) -> None:
    title_box = slide.shapes.add_textbox(
        tokens.CONTENT_LEFT,
        tokens.TITLE_TOP,
        tokens.CONTENT_WIDTH,
        _emu_sum(tokens.TITLE_SIZE, tokens.TITLE_SIZE, tokens.TITLE_SIZE),
    )
    frame = title_box.text_frame
    frame.word_wrap = True  # claims wrap to two lines rather than overflowing the slide
    para = frame.paragraphs[0]
    para.text = action_title
    para.alignment = PP_ALIGN.LEFT
    _style_run(para.runs[0], tokens.TITLE_SIZE, tokens.INK, bold=True)


def _add_accent_rule(slide: Slide) -> None:
    """Draw the single amber rule that separates title from body on every slide."""
    rule_top = Emu(int(tokens.CONTENT_TOP) - int(Inches(0.25)))
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,  # a thin filled bar used as a hairline accent rule
        tokens.CONTENT_LEFT,
        rule_top,
        Inches(1.6),
        Pt(3),
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = tokens.ACCENT_RULE
    rule.line.fill.background()  # no border on the rule — it is a flat accent bar
    rule.shadow.inherit = False  # kill the default drop-shadow (anti-slop)


def _add_bullets(slide: Slide, bullets: tuple[str, ...]) -> None:
    body_height = Emu(int(tokens.SLIDE_HEIGHT) - int(tokens.CONTENT_TOP) - int(tokens.CONTENT_LEFT))
    body_box = slide.shapes.add_textbox(
        tokens.CONTENT_LEFT,
        tokens.CONTENT_TOP,
        tokens.CONTENT_WIDTH,
        body_height,
    )
    frame = body_box.text_frame
    frame.word_wrap = True
    for index, bullet in enumerate(bullets):
        # The first paragraph already exists on a fresh text frame; subsequent
        # bullets get their own paragraph so spacing-after applies per bullet.
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.text = f"•  {bullet}"  # explicit bullet glyph (we own the layout)
        para.space_after = tokens.BULLET_GAP
        _style_run(para.runs[0], tokens.BODY_SIZE, tokens.INK, bold=False)


def _style_run(run: _Run, size: Length, colour: RGBColor, *, bold: bool) -> None:
    """Apply the template font/size/colour to a text run (single styling path)."""
    font = run.font
    font.name = tokens.BODY_FONT
    font.size = size
    font.bold = bold
    font.color.rgb = colour
